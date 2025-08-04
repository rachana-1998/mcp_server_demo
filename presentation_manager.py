import json
from pptx import Presentation
from pptx.util import Inches

def generate_slide(json_file_path, output_pptx_path="/home/user/MyProject/Presentation/json_generated_presentation.pptx"):
    print("i am in presentation_manager file")
    try:
        with open(json_file_path, "r") as f:
            slides_data = json.load(f)
    except FileNotFoundError:
        print(f"Error: JSON file '{json_file_path}' not found.")
        return
    except json.JSONDecodeError:
        print("Error: Invalid JSON format.")
        return

    # Get slides list from JSON
    slides = slides_data.get("slides", [])

    # Create a new presentation
    presentation = Presentation()

    for slide_info in slides:
        slide_type = slide_info.get("slide_type")

        if slide_type == "title":
            slide_layout = presentation.slide_layouts[0]  # Title Slide
            slide = presentation.slides.add_slide(slide_layout)
            slide.shapes.title.text = slide_info.get("title", "")
            subtitle = slide_info.get("subtitle", "")
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = subtitle

        elif slide_type == "bullet":
            slide_layout = presentation.slide_layouts[1]  # Title and Content
            slide = presentation.slides.add_slide(slide_layout)
            slide.shapes.title.text = slide_info.get("title", "")

            bullets = slide_info.get("bullets", [])
            content = slide.placeholders[1]
            tf = content.text_frame
            tf.clear()

            for bullet in bullets:
                p = tf.add_paragraph()
                p.text = bullet

        elif slide_type == "image":
            slide_layout = presentation.slide_layouts[5]  # Title Only
            slide = presentation.slides.add_slide(slide_layout)
            slide.shapes.title.text = slide_info.get("title", "")

            image_path = slide_info.get("image_path")
            if image_path:
                try:
                    slide.shapes.add_picture(image_path, Inches(1), Inches(2), width=Inches(6))
                except Exception as e:
                    print(f"Failed to load image '{image_path}': {e}")
            else:
                print("No image_path provided for image slide")

        else:
            print(f"Unknown slide type: {slide_type}")

    # Save the presentation
    presentation.save(output_pptx_path)
    print(f"✅ Presentation saved as '{output_pptx_path}'")
